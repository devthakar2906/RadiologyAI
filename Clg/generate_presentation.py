from pathlib import Path

from pptx import Presentation


BASE = Path(r"E:\Dev Microsoft\ChatGPT\RadiologyAI\Clg")
SOURCE = BASE / "Sample Presentation.pptx"
OUTPUT = BASE / "RadiologyAI_Final_Presentation.pptx"


def set_bullets(text_frame, lines):
    text_frame.clear()
    for idx, line in enumerate(lines):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = line
        p.level = 0


def add_title_content(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    set_bullets(body, bullets)
    return slide


def add_section(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def main():
    prs = Presentation(str(SOURCE))

    # Title slide
    slide1 = prs.slides[0]
    slide1.shapes.title.text = "Radiology AI Platform"
    subtitle_shapes = [s for s in slide1.shapes if hasattr(s, "text") and s.name == "Subtitle 2"]
    if len(subtitle_shapes) >= 2:
        subtitle_shapes[0].text = (
            "Prepared by\n"
            "Name - Dev Thakar\n"
            "Enrollment No - 230673107042\n"
            "Guided By - Komal Thummar"
        )
        subtitle_shapes[1].text = (
            "Computer Engineering\n"
            "SAL INSTITUTE OF TECHNOLOGY AND ENGINEERING RESEARCH"
        )

    # Table of contents
    slide2 = prs.slides[1]
    slide2.shapes.title.text = "Table of Content"
    toc = [
        "Abstract",
        "Company Profile",
        "Project Overview & Objectives",
        "Software Requirements Specification",
        "Tools & Technologies Used",
        "Key Functionalities / Features",
        "System Workflow & Architecture",
        "Progress Overview",
        "Implementation Highlights",
        "Results & Outcomes",
        "Working Demo",
        "Conclusion & Future Scope",
        "References",
    ]
    set_bullets(slide2.placeholders[1].text_frame, toc)

    add_title_content(
        prs,
        "Abstract",
        [
            "Radiology AI Platform is a full-stack system that helps convert radiology findings into structured reports.",
            "It supports live dictation, uploaded audio refinement, report generation from findings, secure storage, and PDF export.",
            "The project is designed to reduce manual reporting effort and improve consistency in report formatting.",
            "It combines FastAPI, React, PostgreSQL, Redis, Celery, and Hugging Face APIs in one workflow.",
        ],
    )

    add_title_content(
        prs,
        "Company Profile",
        [
            "Internship Organization: FICE Education Pvt. Ltd.",
            "Internship Mode: Remote",
            "Industry Mentor: Anjali Singh",
            "Internship Duration: 21/01/2026 to 01/04/2026 (12 Weeks)",
            "Work focused on practical software development, API integration, UI development, and reporting workflow design.",
        ],
    )

    add_title_content(
        prs,
        "Project Overview & Objectives",
        [
            "Project Title: Radiology AI Platform",
            "Purpose: Convert typed or spoken radiology findings into structured reports.",
            "Objectives include transcript refinement, report generation, secure storage, search, edit, and PDF export.",
            "The system is intended to improve speed, organization, and usability in radiology documentation.",
        ],
    )

    add_title_content(
        prs,
        "Software Requirements Specification",
        [
            "Frontend: React, Vite, Tailwind CSS",
            "Backend: FastAPI, SQLAlchemy",
            "Database: PostgreSQL",
            "Queue & Cache: Redis",
            "Background Processing: Celery",
            "AI Integration: Hugging Face APIs",
            "Others: JWT / cookie auth, AES encryption, Alembic, Flower",
        ],
    )

    add_title_content(
        prs,
        "Tools & Technologies Used",
        [
            "Python for backend implementation",
            "JavaScript and React for frontend development",
            "PostgreSQL for persistent report storage",
            "Redis and Celery for asynchronous processing",
            "Quill editor and jsPDF for editing and export",
            "Git-based iterative development and testing",
        ],
    )

    add_title_content(
        prs,
        "Key Functionalities / Features",
        [
            "Doctor and admin authentication",
            "Live browser dictation and uploaded audio support",
            "AI-based transcript refinement",
            "Structured report generation from findings",
            "Editable structured report panel",
            "Encrypted report storage",
            "Search, filter, save, and PDF export",
            "Concurrent background processing with Celery and Redis",
        ],
    )

    add_section(prs, "System Workflow & Architecture", "Overview of how the platform processes findings into reports")
    add_title_content(
        prs,
        "System Workflow & Architecture",
        [
            "User logs in as doctor or admin.",
            "Findings are entered through typing, live dictation, or uploaded audio.",
            "Uploaded audio is refined through API-based processing.",
            "Findings are converted into structured report sections.",
            "Report can be edited, saved, searched, and downloaded.",
            "Heavy tasks are handled through Celery workers and Redis.",
        ],
    )

    add_title_content(
        prs,
        "Progress Overview",
        [
            "Completed: authentication, dashboard, speech input, structured report generation, save flow, search, PDF export",
            "Completed: encrypted storage, admin filtering, async processing, Flower monitoring",
            "Ongoing improvements: report quality refinement, presentation assets, documentation polishing",
            "Project status: functional prototype with end-to-end working flow",
        ],
    )

    add_title_content(
        prs,
        "Implementation Highlights",
        [
            "Frontend screens include landing page, login/signup, and report dashboard.",
            "Backend services include auth APIs, report APIs, template service, AI service, and worker tasks.",
            "Database stores encrypted report data and user-specific history.",
            "Structured reports are generated through a template-driven AI-assisted flow.",
            "PDF export provides a clean downloadable report for final sharing.",
        ],
    )

    add_title_content(
        prs,
        "Results & Outcomes",
        [
            "System successfully converts findings into structured report output.",
            "Report history can be saved and searched.",
            "Doctors and admins have different access levels.",
            "Audio and findings processing can run asynchronously.",
            "The platform demonstrates practical use of AI in healthcare-oriented documentation.",
        ],
    )

    add_title_content(
        prs,
        "Working Demo",
        [
            "Demo Flow:",
            "1. Login to the platform",
            "2. Enter or dictate findings",
            "3. Refine transcription from uploaded audio",
            "4. Generate structured report",
            "5. Edit and save report",
            "6. Search or export as PDF",
        ],
    )

    add_title_content(
        prs,
        "Conclusion & Future Scope",
        [
            "The project provides a structured and practical reporting workflow for radiology documentation.",
            "It reduces repetitive effort and improves organization of report output.",
            "Future scope includes more study-specific templates, multilingual support, and stronger system integration.",
            "The project can be extended further for hospital-grade deployment and broader clinical use.",
        ],
    )

    add_title_content(
        prs,
        "References",
        [
            "FastAPI Documentation",
            "React Documentation",
            "PostgreSQL Documentation",
            "Redis Documentation",
            "Celery Documentation",
            "Hugging Face Documentation",
            "RadReport",
            "GTU Project Report Guidelines",
        ],
    )

    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()
